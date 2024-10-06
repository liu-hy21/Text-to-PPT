import json
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt


def get_time():
    return datetime.datetime.now().strftime('%Y-%m-%d_%H:%M:%S')

def set_font(p, size, name='Microsoft YaHei'):
    p.font.size = Pt(size)
    p.font.name = name


def add_textbox(slide, content, first_textbox, left_offset, top_offset, font_size=20):
    left = Inches(left_offset)
    top = Inches(top_offset)
    width = Inches(5.5)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = content['text']

    set_font(text_frame.paragraphs[0], font_size)  # 设置字体
    return first_textbox


def add_image(slide, content, left_offset, top_offset):
    """添加图片到幻灯片"""
    img_path = content['path']
    left = Inches(left_offset)
    top = Inches(top_offset)

    slide.shapes.add_picture(img_path, left, top, width=Inches(5), height=Inches(3))


def add_table(slide, content, left_offset, top_offset):
    """解析表格并添加到幻灯片"""
    table_data = content['path'].split('\n')
    headers = table_data[0].split('|')[1:-1]  # 去除两侧的分隔符
    rows = [row.split('|')[1:-1] for row in table_data[2:]]  # 处理表格行数据

    # 添加表格到幻灯片
    left = Inches(left_offset)
    top = Inches(top_offset)
    rows_count = len(rows)
    cols_count = len(headers)

    table_shape = slide.shapes.add_table(rows_count + 1, cols_count, left, top, Inches(6), Inches(2)).table

    # 设置表头
    for col_index, header in enumerate(headers):
        table_shape.cell(0, col_index).text = header.strip()

    # 填充表格数据
    for row_index, row in enumerate(rows, start=1):
        for col_index, cell_text in enumerate(row):
            table_shape.cell(row_index, col_index).text = cell_text.strip()


def handle_content(slide, content, is_first_textbox_exist):
    if content['type'] == 'textbox':
        # 根据是否是第一个文本框来设置不同的位置
        if is_first_textbox_exist == False:
            is_first_textbox_exist = True
            left_offset = 0.75
        else:
            left_offset = 6.75
        is_first_textbox_exist = add_textbox(slide, content, is_first_textbox_exist, left_offset, 1.5)
    elif content['type'] == 'image':
        # 添加图片并根据是否是第一个文本框设置不同的位置
        left_offset = 1 if not is_first_textbox_exist else 6
        add_image(slide, content, left_offset, 2)
    elif content['type'] == 'table':
        # 处理表格并根据是否是第一个文本框设置不同的位置
        left_offset = 1 if not is_first_textbox_exist else 6
        add_table(slide, content, left_offset, 2)

    return is_first_textbox_exist


def generate_pptx(data):
    prs = Presentation('./template.pptx')
    # 逐页处理JSON中的内容
    for page in data:
        # 封面页
        if page['pageNo'] == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            left = Inches(7.2)
            top = Inches(3)
            width = Inches(5.8)
            height = Inches(2)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = page['pageTitle']
            set_font(text_frame.paragraphs[0], 50)
            continue
        # 目录页
        elif page['pageNo'] == 1:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            left = Inches(6)
            top = Inches(1)
            width = Inches(4)
            height = Inches(4.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            if isinstance(page['pageContent'][0]['text'], list):
                text_frame.text = '\n'.join(page['pageContent'][0]['text'])  # 将列表元素拼接为一个字符串
            else:
                text_frame.text = page['pageContent'][0]['text']  # 如果不是列表，直接使用
            for i in range(len(text_frame.paragraphs)):
                set_font(text_frame.paragraphs[i], 26)
            continue
        else:
            slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(slide_layout)
            is_first_textbox_exist = False
            for content in page['pageContent']:
                is_first_textbox_exist = handle_content(slide, content, is_first_textbox_exist)

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    add_textbox(slide, {'text': 'Thank you !'}, False, 6, 2, 50)

    timestamp = get_time()
    # 保存生成的PPT
    prs.save(f'./output_prs_{timestamp}.pptx')


if __name__ == '__main__':
    # 读取JSON文件
    with open('input.json', 'r') as f:
        data = json.load(f)

    # 生成PPTX文件
    generate_pptx(data)
